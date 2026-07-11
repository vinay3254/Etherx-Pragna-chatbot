"""Markdown-outline based content generation shared by all document formats
(Word, Excel, PDF, PowerPoint). See docs/superpowers/specs/2026-07-11-ai-document-generation-design.md.
"""
import re

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table as PdfTable, TableStyle

from services.llm import generate_completion

_TABLE_SEPARATOR_RE = re.compile(r'^\|?\s*:?-{2,}:?\s*(\|\s*:?-{2,}:?\s*)*\|?$')


def _parse_markdown_outline(text):
    """Parse a markdown outline into {"title": str, "sections": [{"heading", "bullets", "table"}]}.

    A section ends up with either a populated `bullets` list (`table: None`) or a
    populated `table` (`bullets: []`) — never both, even if the source text mixed them.
    """
    title = ""
    sections = []
    current = None

    for raw_line in (text or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue

        if line.startswith('# '):
            title = line[2:].strip()
            continue

        if line.startswith('## '):
            current = {"heading": line[3:].strip(), "bullets": [], "table": []}
            sections.append(current)
            continue

        if current is None:
            continue

        if line.startswith('|'):
            if _TABLE_SEPARATOR_RE.match(line):
                continue
            cells = [c.strip() for c in line.strip('|').split('|')]
            current["table"].append(cells)
            continue

        if line.startswith('- ') or line.startswith('* '):
            current["bullets"].append(line[2:].strip())
            continue

    for section in sections:
        if section["table"]:
            section["bullets"] = []
        else:
            section["table"] = None

    return {"title": title or "Untitled Document", "sections": sections}


def generate_document_structure(prompt, language="en"):
    """Ask the LLM for a markdown outline about `prompt` and parse it into a document structure."""
    system_prompt = (
        "You write structured outlines for documents. Given a subject, respond ONLY with "
        "a markdown outline in this exact shape:\n"
        "# <Title>\n"
        "## <Section heading>\n"
        "- <bullet point>\n"
        "- <bullet point>\n"
        "## <Another section heading>\n"
        "| <column> | <column> |\n"
        "| <value> | <value> |\n\n"
        "Use bullet points for narrative sections and a markdown table only for sections "
        "that are genuinely tabular data. Include 3-6 sections. Do not include any text "
        "outside the outline."
    )
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": prompt},
    ]
    outline_text = generate_completion(messages, language=language)
    return _parse_markdown_outline(outline_text)


def _build_docx(structure, filepath):
    doc = Document()
    doc.add_heading(structure["title"], level=0)
    for section in structure["sections"]:
        doc.add_heading(section["heading"], level=1)
        if section["table"]:
            rows = section["table"]
            table = doc.add_table(rows=len(rows), cols=len(rows[0]))
            table.style = "Table Grid"
            for r, row in enumerate(rows):
                for c, cell in enumerate(row):
                    table.cell(r, c).text = cell
        else:
            for bullet in section["bullets"]:
                doc.add_paragraph(bullet, style="List Bullet")
    doc.save(filepath)


def _build_pdf(structure, filepath):
    doc = SimpleDocTemplate(filepath, pagesize=letter)
    styles = getSampleStyleSheet()
    story = [Paragraph(structure["title"], styles["Title"]), Spacer(1, 12)]
    for section in structure["sections"]:
        story.append(Paragraph(section["heading"], styles["Heading2"]))
        story.append(Spacer(1, 6))
        if section["table"]:
            pdf_table = PdfTable(section["table"])
            pdf_table.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
            ]))
            story.append(pdf_table)
        else:
            for bullet in section["bullets"]:
                story.append(Paragraph(f"&bull; {bullet}", styles["Normal"]))
        story.append(Spacer(1, 12))
    doc.build(story)


def _build_pptx(structure, filepath):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = structure["title"]

    for section in structure["sections"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section["heading"]
        body = slide.placeholders[1].text_frame
        body.clear()
        lines = section["bullets"] if section["bullets"] else [
            " | ".join(row) for row in (section["table"] or [])
        ]
        for i, line in enumerate(lines):
            if i == 0:
                body.text = line
            else:
                body.add_paragraph().text = line
    prs.save(filepath)


def _sanitize_sheet_name(name):
    cleaned = re.sub(r'[\[\]:*?/\\]', '', name or '').strip()
    return cleaned[:31] or "Sheet1"


def _build_xlsx(structure, filepath):
    wb = Workbook()
    wb.remove(wb.active)

    table_section = next((s for s in structure["sections"] if s["table"]), None)
    if table_section:
        ws = wb.create_sheet(_sanitize_sheet_name(table_section["heading"]))
        for row in table_section["table"]:
            ws.append(row)
    else:
        ws = wb.create_sheet("Summary")
        ws.append(["Section", "Content"])
        for section in structure["sections"]:
            ws.append([section["heading"], "; ".join(section["bullets"])])

    wb.save(filepath)
